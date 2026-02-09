"""
FastAPI Server for BI Search Agent.

Provides REST and streaming endpoints for the LangGraph agent.
"""

import os
import json
import asyncio
import logging
from typing import Optional
from contextlib import asynccontextmanager

import uuid as _uuid

from fastapi import FastAPI, HTTPException, Request, Query, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, FileResponse, RedirectResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field
from pathlib import Path
from sse_starlette.sse import EventSourceResponse

from agent import (
    run_search_agent,
    run_search_agent_stream,
    get_mcp_client,
    set_request_jwt_token,
    reset_request_jwt_token,
    LLMClientSelector,
)

from database import (
    init_db,
    get_all_dashboards,
    get_dashboard,
    get_dashboard_by_share_id,
    create_dashboard,
    update_dashboard,
    delete_dashboard,
    publish_dashboard,
    unpublish_dashboard,
    save_uploaded_image,
    get_uploaded_image,
    DashboardSchema,
    DashboardCreateSchema,
    DashboardUpdateSchema,
    PublishResponseSchema,
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)


@asynccontextmanager
async def lifespan(app: FastAPI):
    """Application lifespan manager."""
    logger.info("Starting BI Search Agent server...")
    # Initialize database
    await init_db()
    logger.info("Database initialized.")
    yield
    # Cleanup
    logger.info("Shutting down...")
    mcp_client = get_mcp_client()
    await mcp_client.close()


app = FastAPI(
    title="BI Search Agent",
    description="AI-powered Business Intelligence search agent using LangGraph and MCP",
    version="1.0.0",
    lifespan=lifespan,
)

# CORS configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Configure appropriately for production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Static files configuration - serve the built React frontend
BASE_DIR = Path(__file__).resolve().parent
DIST_DIR = BASE_DIR.parent / "dist"  # Frontend build output

# Mount static assets (JS, CSS, etc.)
if DIST_DIR.exists():
    app.mount("/assets", StaticFiles(directory=DIST_DIR / "assets"), name="assets")


@app.get("/")
async def root():
    """Serve the React app."""
    index_file = DIST_DIR / "index.html"
    if index_file.exists():
        return FileResponse(index_file)
    return {"message": "Frontend not built. Run 'npm run build' in the project root."}


# Request/Response Models

class SearchRequest(BaseModel):
    """Search request body."""
    query: str = Field(..., description="The search query")
    conversation_id: Optional[str] = Field(None, description="Conversation ID for context")
    conversation_history: Optional[list[dict]] = Field(default_factory=list, description="Previous conversation turns")
    llm_provider: str = Field("ollama", description="LLM provider (anthropic or ollama)")
    llm_model: Optional[str] = Field(None, description="Specific model to use")
    enabled_tools: Optional[list[str]] = Field(None, description="List of enabled tool names")
    stream: bool = Field(True, description="Whether to stream the response")


class SearchResponse(BaseModel):
    """Non-streaming search response."""
    query: str
    response: str
    sources: list[dict] = []
    chart_configs: list[dict] = []
    thinking_steps: list[dict] = []
    error: Optional[str] = None


class ToolsResponse(BaseModel):
    """Available tools response."""
    tools: list[dict]
    count: int


class ModelsResponse(BaseModel):
    """Available models response."""
    providers: dict[str, list[str]]
    defaults: dict[str, str]


class HealthResponse(BaseModel):
    """Health check response."""
    status: str
    mcp_session_stats: dict


class FieldInfo(BaseModel):
    """Field information from MCP tool."""
    name: str
    type: str  # keyword, date, numeric, derived
    description: Optional[str] = None


class DataSourceInfo(BaseModel):
    """Data source (MCP tool) information."""
    id: str
    name: str
    description: Optional[str] = None
    fields: list[FieldInfo]
    date_fields: list[str]
    groupable_fields: list[str]  # Fields that can be used for group_by


class DataSourcesResponse(BaseModel):
    """Available data sources response."""
    sources: list[DataSourceInfo]


class ChartDataRequest(BaseModel):
    """Chart data request body."""
    index: str = Field(..., description="OpenSearch index name or tool name")
    x_field: str = Field(..., description="Field for X-axis (grouping)")
    y_field: Optional[str] = Field(None, description="Field for Y-axis (metric)")
    series_field: Optional[str] = Field(None, description="Field to split data into multiple series")
    aggregation: str = Field("count", description="Aggregation type (count, sum, avg, min, max)")
    chart_type: str = Field("bar", description="Chart type")
    filters: list[dict] = Field(default_factory=list, description="Filters to apply")


class ChartDataResponse(BaseModel):
    """Chart data response."""
    labels: list[str]
    datasets: list[dict]
    error: Optional[str] = None  # Error message if data fetch failed


# Endpoints

@app.get("/health", response_model=HealthResponse)
async def health_check():
    """Health check endpoint."""
    mcp_client = get_mcp_client()
    return HealthResponse(
        status="healthy",
        mcp_session_stats=mcp_client.get_session_stats(),
    )


def parse_tool_fields(description: str) -> tuple[list[FieldInfo], list[str], list[str]]:
    """
    Parse field information from MCP tool description.

    Returns: (fields, date_fields, groupable_fields)
    """
    import re

    fields = []
    date_fields = []
    groupable_fields = []

    # Parse <fields> section
    fields_match = re.search(r'<fields>(.*?)</fields>', description, re.DOTALL)
    if fields_match:
        fields_text = fields_match.group(1)

        # Parse keyword fields
        keyword_match = re.search(r'keyword:\s*([^\n]+)', fields_text)
        if keyword_match:
            keyword_fields = [f.strip() for f in keyword_match.group(1).split(',')]
            for field in keyword_fields:
                if field:
                    fields.append(FieldInfo(name=field, type='keyword'))
                    groupable_fields.append(field)

        # Parse date fields
        date_match = re.search(r'date:\s*([^\n]+)', fields_text)
        if date_match:
            date_field_names = [f.strip() for f in date_match.group(1).split(',')]
            for field in date_field_names:
                if field:
                    fields.append(FieldInfo(name=field, type='date'))
                    date_fields.append(field)

        # Parse derived fields (like year)
        derived_match = re.search(r'year:\s*([^\n]+)', fields_text)
        if derived_match:
            fields.append(FieldInfo(name='year', type='derived', description='Year derived from date field'))
            groupable_fields.append('year')

    # Parse <field_context> for descriptions
    context_match = re.search(r'<field_context>(.*?)</field_context>', description, re.DOTALL)
    if context_match:
        context_text = context_match.group(1)

        # Extract field descriptions
        for field in fields:
            # Look for "field_name: description" pattern
            desc_match = re.search(rf'{field.name}:\s*([^\n]+)', context_text)
            if desc_match:
                field.description = desc_match.group(1).strip()

    return fields, date_fields, groupable_fields


@app.get("/data-sources", response_model=DataSourcesResponse)
async def get_data_sources():
    """
    Get available data sources from MCP tools with their field metadata.

    Parses tool descriptions to extract field information for chart creation.
    """
    try:
        mcp_client = get_mcp_client()
        tools = await mcp_client.get_available_tools()

        sources = []
        for tool in tools:
            name = tool.get('name', '')
            description = tool.get('description', '')

            # Only include analytics tools
            if 'analyze' not in name.lower():
                continue

            # Parse fields from description
            fields, date_fields, groupable_fields = parse_tool_fields(description)

            # Create a friendly display name
            display_name = name.replace('_', ' ').title()
            if 'conclusion' in name.lower():
                display_name = 'Events (by Conclusion Date)'
            elif 'all_events' in name.lower():
                display_name = 'All Events'

            sources.append(DataSourceInfo(
                id=name,
                name=display_name,
                description=description[:200] + '...' if len(description) > 200 else description,
                fields=fields,
                date_fields=date_fields,
                groupable_fields=groupable_fields,
            ))

        # If no MCP tools available, raise an error
        if not sources:
            raise HTTPException(status_code=503, detail="No analytics tools available from MCP gateway")

        return DataSourcesResponse(sources=sources)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error fetching data sources: {e}")
        raise HTTPException(status_code=503, detail=f"Failed to fetch data sources from MCP: {str(e)}")


@app.get("/debug/mcp-test")
async def debug_mcp_test():
    """Debug endpoint to test MCP tool call directly."""
    try:
        mcp_client = get_mcp_client()
        result = await mcp_client.call_tool(
            tool_name="analyze_all_events",
            arguments={"group_by": "country", "top_n": 5},
            user_email="anonymous"
        )
        return {"status": "success", "result": result}
    except Exception as e:
        return {"status": "error", "error": str(e)}


@app.post("/chart-data", response_model=ChartDataResponse)
async def get_chart_data(request: Request, body: ChartDataRequest):
    """
    Fetch chart data via MCP gateway.

    Uses the analytics tools (analyze_all_events, analyze_events_by_conclusion)
    to query data with group_by aggregations.

    If series_field is specified, creates multiple datasets by splitting data
    on the series_field values.
    """
    # Set JWT token from header if present
    auth_header = request.headers.get("Authorization", "")
    if auth_header.startswith("Bearer "):
        token = auth_header[7:]
        reset_token = set_request_jwt_token(token)
    else:
        reset_token = None

    try:
        mcp_client = get_mcp_client()

        # Map index/dataSource to tool name
        tool_name = body.index
        if tool_name.startswith("analyze_"):
            pass
        elif "conclusion" in tool_name.lower():
            tool_name = "analyze_events_by_conclusion"
        else:
            tool_name = "analyze_all_events"

        # Build base filter dict and range filter dict from operator
        base_filters = {}
        range_filters = {}
        if body.filters:
            for f in body.filters:
                field = f.get("field")
                value = f.get("value")
                operator = f.get("operator", "eq")
                if not field or value is None:
                    continue
                if operator in ("eq", "in", "contains"):
                    base_filters[field] = value
                elif operator == "neq":
                    # MCP doesn't support neq natively — skip for now
                    continue
                elif operator in ("gt", "gte", "lt", "lte"):
                    if field not in range_filters:
                        range_filters[field] = {}
                    range_filters[field][operator] = value

        # If series_field is specified, we need to split data by that field
        if body.series_field:
            return await fetch_chart_data_with_series(
                mcp_client, tool_name, body, base_filters, range_filters
            )

        # Standard single-series fetch
        tool_arguments = {
            "group_by": body.x_field,
            "top_n": 20,
        }
        if base_filters:
            tool_arguments["filters"] = json.dumps(base_filters)
        if range_filters:
            tool_arguments["range_filters"] = json.dumps(range_filters)

        logger.info(f"Calling MCP tool: {tool_name} with args: {tool_arguments}")

        result = await mcp_client.call_tool(
            tool_name=tool_name,
            arguments=tool_arguments,
            user_email="anonymous"
        )

        logger.info(f"MCP tool result: {json.dumps(result)[:500]}")

        if "error" in result:
            error_msg = result.get("error", {}).get("message", str(result.get("error")))
            logger.warning(f"MCP tool call failed: {error_msg}")
            return ChartDataResponse(labels=[], datasets=[], error=f"MCP tool error: {error_msg}")

        # Check for structured error from MCP
        structured = result.get("structuredContent", {})
        if structured.get("error"):
            error_msg = structured["error"]
            logger.warning(f"MCP structured error: {error_msg}")
            return ChartDataResponse(labels=[], datasets=[], error=f"Data source error: {error_msg}")

        # Parse structuredContent (preferred) or content
        if structured and structured.get("status") == "success":
            return parse_structured_content(structured, body)

        # Fallback: Parse the result from MCP content
        content = result.get("content", [])
        if not content:
            return ChartDataResponse(labels=[], datasets=[], error="No data returned from MCP tool")

        text_content = ""
        for item in content:
            if isinstance(item, dict) and item.get("type") == "text":
                text_content = item.get("text", "")
                break
            elif isinstance(item, str):
                text_content = item
                break

        if not text_content:
            return ChartDataResponse(labels=[], datasets=[], error="Empty response from MCP tool")

        try:
            analytics_response = json.loads(text_content)
            return parse_analytics_response(analytics_response, body)
        except json.JSONDecodeError:
            logger.error(f"Failed to parse analytics response: {text_content[:500]}")
            return ChartDataResponse(labels=[], datasets=[], error="Failed to parse MCP response")

    except Exception as e:
        logger.error(f"Error fetching chart data: {e}")
        return ChartDataResponse(labels=[], datasets=[], error=f"Error: {str(e)}")

    finally:
        if reset_token:
            reset_request_jwt_token(reset_token)


async def fetch_chart_data_with_series(mcp_client, tool_name: str, body: ChartDataRequest, base_filters: dict, range_filters: dict | None = None) -> ChartDataResponse:
    """
    Fetch chart data split by series_field.

    1. First get unique values of the series_field
    2. For each unique value, query with it as a filter
    3. Combine all results into multiple datasets
    """
    if range_filters is None:
        range_filters = {}

    # Step 1: Get unique values of series_field
    series_args = {
        "group_by": body.series_field,
        "top_n": 10,  # Limit to top 10 series for performance
    }
    if base_filters:
        series_args["filters"] = json.dumps(base_filters)
    if range_filters:
        series_args["range_filters"] = json.dumps(range_filters)

    logger.info(f"Fetching series values with: {series_args}")

    series_result = await mcp_client.call_tool(
        tool_name=tool_name,
        arguments=series_args,
        user_email="anonymous"
    )

    # Parse series values from result
    series_values = []
    structured = series_result.get("structuredContent", {})
    if structured and structured.get("status") == "success":
        buckets = structured.get("aggregations", {}).get("group_by", {}).get("buckets", [])
        series_values = [str(b.get("key", "")) for b in buckets if b.get("key")]
    else:
        # Try parsing from content
        content = series_result.get("content", [])
        for item in content:
            if isinstance(item, dict) and item.get("type") == "text":
                try:
                    data = json.loads(item.get("text", ""))
                    for key, value in data.get("aggregations", {}).items():
                        if isinstance(value, dict) and "buckets" in value:
                            series_values = [str(b.get("key", "")) for b in value.get("buckets", []) if b.get("key")]
                            break
                except:
                    pass
                break

    if not series_values:
        logger.warning(f"No series values found for field: {body.series_field}")
        return ChartDataResponse(labels=[], datasets=[], error=f"No values found for series field: {body.series_field}")

    logger.info(f"Found series values: {series_values}")

    # Step 2: For each series value, get the data grouped by x_field
    all_labels = set()
    series_data = {}  # {series_value: {x_value: count}}

    for series_val in series_values:
        # Build filters with series value
        series_filters = {**base_filters, body.series_field: series_val}

        tool_args = {
            "group_by": body.x_field,
            "top_n": 20,
            "filters": json.dumps(series_filters),
        }
        if range_filters:
            tool_args["range_filters"] = json.dumps(range_filters)

        logger.info(f"Fetching data for series '{series_val}': {tool_args}")

        result = await mcp_client.call_tool(
            tool_name=tool_name,
            arguments=tool_args,
            user_email="anonymous"
        )

        # Parse the result
        buckets = []
        structured = result.get("structuredContent", {})
        if structured and structured.get("status") == "success":
            buckets = structured.get("aggregations", {}).get("group_by", {}).get("buckets", [])
        else:
            content = result.get("content", [])
            for item in content:
                if isinstance(item, dict) and item.get("type") == "text":
                    try:
                        data = json.loads(item.get("text", ""))
                        for key, value in data.get("aggregations", {}).items():
                            if isinstance(value, dict) and "buckets" in value:
                                buckets = value.get("buckets", [])
                                break
                    except:
                        pass
                    break

        # Store data for this series
        series_data[series_val] = {}
        for bucket in buckets:
            x_val = str(bucket.get("key", ""))
            count = bucket.get("doc_count", bucket.get("count", 0))
            if x_val:
                all_labels.add(x_val)
                series_data[series_val][x_val] = count

    # Step 3: Build the response with aligned data
    labels = sorted(list(all_labels))
    datasets = []

    for series_val in series_values:
        data_dict = series_data.get(series_val, {})
        # Align data with labels (fill missing with 0)
        data = [data_dict.get(label, 0) for label in labels]

        if body.chart_type in ("pie", "funnel"):
            datasets.append({
                "name": series_val,
                "data": [{"name": label, "value": val} for label, val in zip(labels, data)]
            })
        else:
            datasets.append({
                "name": series_val,
                "data": data
            })

    logger.info(f"Built {len(datasets)} series with {len(labels)} labels")

    return ChartDataResponse(labels=labels, datasets=datasets)


def parse_structured_content(structured: dict, body: ChartDataRequest) -> ChartDataResponse:
    """Parse MCP structuredContent format into chart data."""
    # Try to get aggregation buckets
    aggregations = structured.get("aggregations", {})
    group_by = aggregations.get("group_by", {})
    buckets = group_by.get("buckets", [])

    if not buckets:
        # Try chart_config as fallback
        chart_config = structured.get("chart_config", [])
        if chart_config and len(chart_config) > 0:
            config = chart_config[0]
            labels = config.get("labels", [])
            data = config.get("data", [])
            if labels and data:
                if body.chart_type == "pie":
                    return ChartDataResponse(
                        labels=labels,
                        datasets=[{
                            "name": config.get("title", f"{body.x_field} distribution"),
                            "data": [{"name": label, "value": value} for label, value in zip(labels, data)]
                        }]
                    )
                return ChartDataResponse(
                    labels=labels,
                    datasets=[{
                        "name": config.get("title", f"Count by {body.x_field}"),
                        "data": data
                    }]
                )
        return ChartDataResponse(labels=[], datasets=[], error="No aggregation data found")

    # Extract labels and values from buckets
    labels = [str(b.get("key", "")) for b in buckets]
    values = [b.get("doc_count", b.get("count", 0)) for b in buckets]

    if body.chart_type == "pie":
        return ChartDataResponse(
            labels=labels,
            datasets=[{
                "name": f"{body.x_field} distribution",
                "data": [{"name": label, "value": value} for label, value in zip(labels, values)]
            }]
        )

    return ChartDataResponse(
        labels=labels,
        datasets=[{
            "name": f"Count by {body.x_field}",
            "data": values
        }]
    )


def parse_analytics_response(response: dict, body: ChartDataRequest) -> ChartDataResponse:
    """Parse analytics tool response into chart data format."""
    # The analytics tools return aggregations in a specific format
    aggregations = response.get("aggregations", {})

    # Try to find the group_by buckets
    buckets = []
    for key, value in aggregations.items():
        if isinstance(value, dict) and "buckets" in value:
            buckets = value.get("buckets", [])
            break

    if not buckets:
        # Maybe it's a flat structure
        if "buckets" in aggregations:
            buckets = aggregations.get("buckets", [])

    if not buckets:
        logger.warning(f"No buckets found in response: {json.dumps(response)[:500]}")
        return ChartDataResponse(labels=[], datasets=[], error="No aggregation data found in response")

    labels = [str(b.get("key", "")) for b in buckets]
    values = [b.get("doc_count", 0) for b in buckets]

    if body.chart_type in ("pie", "funnel"):
        return ChartDataResponse(
            labels=labels,
            datasets=[{
                "name": f"{body.x_field} distribution",
                "data": [{"name": label, "value": value} for label, value in zip(labels, values)]
            }]
        )

    return ChartDataResponse(
        labels=labels,
        datasets=[{
            "name": f"Count by {body.x_field}",
            "data": values
        }]
    )


def parse_opensearch_response(response: dict, body: ChartDataRequest) -> ChartDataResponse:
    """Parse OpenSearch aggregation response into chart data format."""
    aggs = response.get("aggregations", {})
    categories = aggs.get("categories", {})
    buckets = categories.get("buckets", [])

    if not buckets:
        return ChartDataResponse(labels=[], datasets=[], error="No data found in aggregation response")

    labels = [str(b.get("key", "")) for b in buckets]

    # For count aggregation, use doc_count
    if body.aggregation == "count":
        values = [b.get("doc_count", 0) for b in buckets]
    else:
        # For other aggregations, use the metric value
        values = [b.get("metric", {}).get("value", 0) for b in buckets]

    if body.chart_type in ("pie", "funnel"):
        return ChartDataResponse(
            labels=labels,
            datasets=[{
                "name": f"{body.x_field} distribution",
                "data": [{"name": label, "value": value} for label, value in zip(labels, values)]
            }]
        )

    return ChartDataResponse(
        labels=labels,
        datasets=[{
            "name": f"{body.aggregation} of {body.y_field or 'records'}",
            "data": values
        }]
    )


@app.get("/tools", response_model=ToolsResponse)
async def get_tools(
    request: Request,
    user_email: str = Query("anonymous", description="User email for tool access"),
):
    """Get available tools from the MCP gateway."""
    # Set JWT token from header if present
    auth_header = request.headers.get("Authorization", "")
    if auth_header.startswith("Bearer "):
        token = auth_header[7:]
        reset_token = set_request_jwt_token(token)
    else:
        reset_token = None

    try:
        mcp_client = get_mcp_client()
        tools = await mcp_client.get_available_tools(user_email)
        return ToolsResponse(tools=tools, count=len(tools))
    except Exception as e:
        logger.warning(f"Failed to get tools from MCP gateway: {e}")
        # Return default tools when MCP is unavailable
        default_tools = [
            {"name": "opensearch_query", "description": "Query OpenSearch for analytics data"},
            {"name": "chart_generator", "description": "Generate chart configurations"},
        ]
        return ToolsResponse(tools=default_tools, count=len(default_tools))
    finally:
        if reset_token:
            reset_request_jwt_token(reset_token)


@app.get("/models", response_model=ModelsResponse)
async def get_models():
    """Get available LLM models."""
    return ModelsResponse(
        providers={
            "anthropic": LLMClientSelector.get_available_models("anthropic"),
            "ollama": LLMClientSelector.get_available_models("ollama"),
        },
        defaults={
            "anthropic": LLMClientSelector.get_default_model("anthropic"),
            "ollama": LLMClientSelector.get_default_model("ollama"),
        }
    )


@app.post("/search")
async def search(request: Request, body: SearchRequest):
    """
    Main search endpoint.

    If stream=True (default), returns Server-Sent Events stream.
    If stream=False, returns JSON response.
    """
    # Set JWT token from header if present
    auth_header = request.headers.get("Authorization", "")
    if auth_header.startswith("Bearer "):
        token = auth_header[7:]
        reset_token = set_request_jwt_token(token)
    else:
        reset_token = None

    # Extract user email from token or use provided
    user_email = "anonymous"  # Could extract from JWT claims

    try:
        if body.stream:
            return EventSourceResponse(
                stream_search(
                    query=body.query,
                    user_email=user_email,
                    conversation_id=body.conversation_id,
                    conversation_history=body.conversation_history,
                    llm_provider=body.llm_provider,
                    llm_model=body.llm_model,
                    enabled_tools=body.enabled_tools,
                    jwt_token=auth_header[7:] if auth_header.startswith("Bearer ") else None,
                ),
                media_type="text/event-stream",
            )
        else:
            # Non-streaming response
            final_state = await run_search_agent(
                query=body.query,
                user_email=user_email,
                conversation_id=body.conversation_id,
                conversation_history=body.conversation_history,
                llm_provider=body.llm_provider,
                llm_model=body.llm_model,
                enabled_tools=body.enabled_tools,
            )

            return SearchResponse(
                query=body.query,
                response=final_state.get("final_response", ""),
                sources=final_state.get("extracted_sources", []),
                chart_configs=final_state.get("chart_configs", []),
                thinking_steps=final_state.get("thinking_steps", []),
                error=final_state.get("error_message"),
            )
    finally:
        if reset_token:
            reset_request_jwt_token(reset_token)


async def stream_search(
    query: str,
    user_email: str,
    conversation_id: Optional[str],
    conversation_history: Optional[list],
    llm_provider: str,
    llm_model: Optional[str],
    enabled_tools: Optional[list],
    jwt_token: Optional[str],
):
    """
    Generator for streaming search events.

    Event types:
    - thinking: Agent thinking/processing step
    - node_start: Node execution started
    - response_start: Final response starting
    - response_char: Single character of response
    - response_end: Final response complete
    - sources: Extracted sources
    - chart_configs: Suggested chart configurations
    - error: Error occurred
    - complete: Agent finished
    """
    # Set JWT context for this stream
    if jwt_token:
        reset_token = set_request_jwt_token(jwt_token)
    else:
        reset_token = None

    try:
        async for event_type, data in run_search_agent_stream(
            query=query,
            user_email=user_email,
            conversation_id=conversation_id,
            conversation_history=conversation_history,
            llm_provider=llm_provider,
            llm_model=llm_model,
            enabled_tools=enabled_tools,
        ):
            if event_type == "node_start":
                yield {
                    "event": "thinking",
                    "data": json.dumps({"type": "node_start", "node": data["node"]}),
                }

            elif event_type == "thinking":
                yield {
                    "event": "thinking",
                    "data": json.dumps({"type": "step", **data}),
                }

            elif event_type == "response_start":
                yield {
                    "event": "response",
                    "data": json.dumps({"type": "start"}),
                }

            elif event_type == "response_char":
                yield {
                    "event": "response",
                    "data": json.dumps({"type": "char", "char": data}),
                }
                # Small delay for typing effect
                await asyncio.sleep(0.002)

            elif event_type == "response_end":
                yield {
                    "event": "response",
                    "data": json.dumps({"type": "end"}),
                }

            elif event_type == "sources":
                yield {
                    "event": "sources",
                    "data": json.dumps(data),
                }

            elif event_type == "chart_configs":
                yield {
                    "event": "charts",
                    "data": json.dumps(data),
                }

            elif event_type == "presentation_config":
                yield {
                    "event": "presentation",
                    "data": json.dumps(data),
                }

            elif event_type == "error":
                yield {
                    "event": "error",
                    "data": json.dumps(data),
                }

            elif event_type == "complete":
                yield {
                    "event": "complete",
                    "data": json.dumps(data),
                }

    except Exception as e:
        logger.error(f"Stream error: {e}")
        yield {
            "event": "error",
            "data": json.dumps({"message": str(e)}),
        }

    finally:
        if reset_token:
            reset_request_jwt_token(reset_token)


@app.post("/chat")
async def chat(request: Request, body: SearchRequest):
    """
    Chat-style endpoint (alias for /search).
    Provided for compatibility with chat interfaces.
    """
    return await search(request, body)


# Dashboard Endpoints

@app.get("/api/dashboards", response_model=list[DashboardSchema])
async def list_dashboards():
    """Get all dashboards."""
    return await get_all_dashboards()


@app.get("/api/dashboards/{dashboard_id}", response_model=DashboardSchema)
async def get_dashboard_by_id(dashboard_id: str):
    """Get a dashboard by ID."""
    dashboard = await get_dashboard(dashboard_id)
    if not dashboard:
        raise HTTPException(status_code=404, detail="Dashboard not found")
    return dashboard


@app.post("/api/dashboards", response_model=DashboardSchema)
async def create_new_dashboard(data: DashboardCreateSchema):
    """Create a new dashboard."""
    return await create_dashboard(data)


@app.put("/api/dashboards/{dashboard_id}", response_model=DashboardSchema)
async def update_existing_dashboard(dashboard_id: str, data: DashboardUpdateSchema):
    """Update an existing dashboard."""
    dashboard = await update_dashboard(dashboard_id, data)
    if not dashboard:
        raise HTTPException(status_code=404, detail="Dashboard not found")
    return dashboard


@app.delete("/api/dashboards/{dashboard_id}")
async def delete_existing_dashboard(dashboard_id: str):
    """Delete a dashboard."""
    success = await delete_dashboard(dashboard_id)
    if not success:
        raise HTTPException(status_code=404, detail="Dashboard not found")
    return {"success": True}


@app.post("/api/dashboards/{dashboard_id}/publish", response_model=PublishResponseSchema)
async def publish_dashboard_endpoint(dashboard_id: str):
    """Publish a dashboard and get a shareable link."""
    result = await publish_dashboard(dashboard_id)
    if not result:
        raise HTTPException(status_code=404, detail="Dashboard not found")
    return result


@app.post("/api/dashboards/{dashboard_id}/unpublish")
async def unpublish_dashboard_endpoint(dashboard_id: str):
    """Unpublish a dashboard."""
    success = await unpublish_dashboard(dashboard_id)
    if not success:
        raise HTTPException(status_code=404, detail="Dashboard not found")
    return {"success": True}


@app.get("/api/shared/{share_id}", response_model=DashboardSchema)
async def get_shared_dashboard(share_id: str):
    """Get a published dashboard by share ID."""
    dashboard = await get_dashboard_by_share_id(share_id)
    if not dashboard:
        raise HTTPException(status_code=404, detail="Dashboard not found or not published")
    return dashboard


@app.post("/api/dashboards/{dashboard_id}/export")
async def export_dashboard(dashboard_id: str):
    """Export a dashboard as JSON."""
    dashboard = await get_dashboard(dashboard_id)
    if not dashboard:
        raise HTTPException(status_code=404, detail="Dashboard not found")

    return {
        "id": dashboard.id,
        "title": dashboard.title,
        "charts": dashboard.charts,
        "layout": dashboard.layout,
        "filters": dashboard.filters,
        "exported_at": datetime.now().isoformat(),
    }


from datetime import datetime


# Optional: WebSocket endpoint for bidirectional streaming
# @app.websocket("/ws/search")
# async def websocket_search(websocket):
#     ...


# Image upload/serve endpoints — stored as BLOBs in the database

ALLOWED_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg"}
MAX_IMAGE_SIZE = 10 * 1024 * 1024  # 10 MB

# Map extensions to MIME types
MIME_TYPES = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".svg": "image/svg+xml",
}


@app.post("/api/images/upload")
async def upload_image(file: UploadFile = File(...)):
    """Upload an image file, store it in the database, and return its URL."""
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided")

    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_IMAGE_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"File type '{ext}' not allowed. Allowed: {', '.join(ALLOWED_IMAGE_EXTENSIONS)}"
        )

    contents = await file.read()
    if len(contents) > MAX_IMAGE_SIZE:
        raise HTTPException(status_code=400, detail="File too large. Maximum size is 10 MB.")

    image_id = _uuid.uuid4().hex[:12]
    content_type = MIME_TYPES.get(ext, file.content_type or "application/octet-stream")

    await save_uploaded_image(
        image_id=image_id,
        filename=file.filename,
        content_type=content_type,
        data=contents,
    )

    return {"url": f"/api/images/{image_id}", "filename": image_id}


@app.get("/api/images/{image_id}")
async def serve_image(image_id: str):
    """Serve an uploaded image from the database."""
    image = await get_uploaded_image(image_id)
    if not image:
        raise HTTPException(status_code=404, detail="Image not found")

    from fastapi.responses import Response
    return Response(
        content=image.data,
        media_type=image.content_type,
        headers={
            "Cache-Control": "public, max-age=31536000, immutable",
        },
    )


# Presentation endpoints

class PresentationExportRequest(BaseModel):
    """Request body for PPTX export."""
    id: str
    title: str
    slides: list[dict]
    theme: Optional[dict] = None


@app.post("/api/presentations/export")
async def export_presentation_pptx(body: PresentationExportRequest):
    """Export a presentation as a PPTX file using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io

        prs = Presentation()
        # Set slide dimensions to 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide_width_emu = prs.slide_width
        slide_height_emu = prs.slide_height

        for slide_data in body.slides:
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            for element in slide_data.get("elements", []):
                el_type = element.get("type")
                x_pct = element.get("x", 0) / 100
                y_pct = element.get("y", 0) / 100
                w_pct = element.get("width", 10) / 100
                h_pct = element.get("height", 10) / 100

                left = int(slide_width_emu * x_pct)
                top = int(slide_height_emu * y_pct)
                width = int(slide_width_emu * w_pct)
                height = int(slide_height_emu * h_pct)

                style = element.get("style", {})

                if el_type == "text":
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = element.get("content", "")

                    # Apply style
                    font_size = style.get("fontSize", 16)
                    p.font.size = Pt(font_size)

                    if style.get("fontWeight") == "bold":
                        p.font.bold = True
                    if style.get("fontStyle") == "italic":
                        p.font.italic = True

                    color_hex = style.get("color", "#000000").lstrip("#")
                    if len(color_hex) == 6:
                        p.font.color.rgb = RGBColor(
                            int(color_hex[0:2], 16),
                            int(color_hex[2:4], 16),
                            int(color_hex[4:6], 16),
                        )

                    align = style.get("textAlign", "left")
                    if align == "center":
                        p.alignment = PP_ALIGN.CENTER
                    elif align == "right":
                        p.alignment = PP_ALIGN.RIGHT

                elif el_type == "shape":
                    shape = slide.shapes.add_shape(
                        1,  # MSO_SHAPE.RECTANGLE
                        left, top, width, height,
                    )
                    bg_hex = style.get("backgroundColor", "#e2e8f0").lstrip("#")
                    if len(bg_hex) == 6:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(
                            int(bg_hex[0:2], 16),
                            int(bg_hex[2:4], 16),
                            int(bg_hex[4:6], 16),
                        )

            # Add speaker notes if present
            notes_text = slide_data.get("notes", "")
            if notes_text:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text

        # Write to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        filename = f"{body.title.replace(' ', '_')}.pptx"

        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"',
            },
        )

    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="python-pptx not installed. Run: pip install python-pptx",
        )
    except Exception as e:
        logger.error(f"PPTX export error: {e}")
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")


# Catch-all route for SPA - MUST be last to not intercept API routes
@app.get("/{path:path}")
async def serve_spa(path: str):
    """Serve static files or fall back to index.html for SPA routing."""
    # Try to serve the exact file
    file_path = DIST_DIR / path
    if file_path.exists() and file_path.is_file():
        return FileResponse(file_path)

    # Fall back to index.html for SPA routing
    index_file = DIST_DIR / "index.html"
    if index_file.exists():
        return FileResponse(index_file)

    raise HTTPException(status_code=404, detail="Not found")


if __name__ == "__main__":
    import uvicorn

    port = int(os.getenv("PORT", "8025"))
    host = os.getenv("HOST", "0.0.0.0")

    uvicorn.run(
        "server:app",
        host=host,
        port=port,
        reload=True,
        log_level="info",
    )
